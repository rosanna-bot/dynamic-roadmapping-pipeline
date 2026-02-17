#!/usr/bin/env python3
"""Dynamic Roadmapping V1 Kick-Off — clean, direct, no fluff."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── Colours ──
DARK   = RGBColor(0x05, 0x00, 0x38)
DEEP   = RGBColor(0x12, 0x0E, 0x4A)
PURPLE = RGBColor(0x42, 0x62, 0xFF)
YELLOW = RGBColor(0xFF, 0xD0, 0x2F)
RED    = RGBColor(0xF2, 0x47, 0x26)
GREEN  = RGBColor(0x00, 0xB8, 0x75)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xFA, 0xFA, 0xFA)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
TXT    = RGBColor(0x1A, 0x1A, 0x2E)
LGRAY  = RGBColor(0xE5, 0xE7, 0xEB)
SOFT   = RGBColor(0x99, 0x99, 0xBB)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BL = prs.slide_layouts[6]

def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def rect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def rrect(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, text, sz=18, b=False, c=TXT, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c; p.font.name = 'Calibri'
    p.alignment = al; p.space_before = Pt(0); p.space_after = Pt(0)
    return tb

def rich(s, l, t, w, h, runs, al=PP_ALIGN.LEFT):
    """runs = [(text, size, bold, color), ...]"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    for text, sz, b, c in runs:
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = 'Calibri'
    return tb

def bullets(s, l, t, w, h, items, sz=15, c=TXT, bc=PURPLE, gap=Pt(8)):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = 'Calibri'
        p.space_before = gap; p.space_after = Pt(2)
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, f'{{{ns}}}pPr')
        for tag in ['buNone']:
            el = pPr.find(f'{{{ns}}}{tag}')
            if el is not None: pPr.remove(el)
        buChar = etree.SubElement(pPr, f'{{{ns}}}buChar'); buChar.set('char', '—')
        buClr = etree.SubElement(pPr, f'{{{ns}}}buClr')
        srgb = etree.SubElement(buClr, f'{{{ns}}}srgbClr')
        srgb.set('val', f'{bc[0]:02X}{bc[1]:02X}{bc[2]:02X}' if isinstance(bc, tuple) else
                  '{:02X}{:02X}{:02X}'.format(*bc))
        buSzPct = etree.SubElement(pPr, f'{{{ns}}}buSzPct'); buSzPct.set('val', '100000')
        pPr.set('indent', str(Emu(Inches(0.2)))); pPr.set('marL', str(Emu(Inches(0.3))))
    return tb

def tbl(s, l, t, w, rows, cw, hbg=DARK):
    n_r, n_c = len(rows), len(rows[0])
    sh = s.shapes.add_table(n_r, n_c, l, t, w, Inches(0.38 * n_r))
    tb = sh.table
    for ci, c in enumerate(cw): tb.columns[ci].width = c
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tb.cell(ri, ci); cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12 if ri > 0 else 11); p.font.name = 'Calibri'
                p.font.bold = ri == 0; p.font.color.rgb = WHITE if ri == 0 else TXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = hbg if ri == 0 else WHITE
    return sh

def stat(s, l, t, num, label, nc=YELLOW, bgc=DEEP):
    rrect(s, l, t, Inches(2.3), Inches(1.3), bgc)
    txt(s, l, t + Inches(0.12), Inches(2.3), Inches(0.6), num, sz=40, b=True, c=nc, al=PP_ALIGN.CENTER)
    txt(s, l + Inches(0.15), t + Inches(0.7), Inches(2), Inches(0.55), label, sz=10, c=SOFT, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# 1. TITLE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
rect(s, Inches(1), Inches(2.65), Inches(0.45), Inches(0.05), YELLOW)
txt(s, Inches(1), Inches(2.0), Inches(4), Inches(0.3), "KICK-OFF · FEB 2026", sz=11, b=True, c=YELLOW)
rich(s, Inches(1), Inches(2.9), Inches(9), Inches(1.8), [
    ("Dynamic\nRoadmapping ", 50, True, WHITE),
    ("V1", 50, True, YELLOW),
])


# ═══════════════════════════════════════════════════
# 2. THE SITUATION
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(1), Inches(10), Inches(1.2),
    "60% of Miro customers already\nuse us for roadmapping.", sz=38, b=True, c=TXT)
txt(s, Inches(1), Inches(2.8), Inches(8), Inches(1),
    "It has the highest frustration score across all agile use cases.\n"
    "Only 34% are very satisfied. $2.7M ARR in customers are\n"
    "explicitly telling us we're not good enough.",
    sz=18, c=MUTED)
txt(s, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
    "We have the demand. We're not meeting it.",
    sz=20, b=True, c=PURPLE)


# ═══════════════════════════════════════════════════
# 3. THE ACTUAL PROBLEM
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.8), Inches(10), Inches(1.2),
    "Roadmaps today are static snapshots\nthat decay the moment they're created.",
    sz=34, b=True, c=WHITE)
txt(s, Inches(1), Inches(2.4), Inches(8), Inches(0.8),
    "Strategy in slides. Feedback in Gong. Execution in Jira.\n"
    "The PM stitches it together manually, 4-6 hours a week.",
    sz=17, c=SOFT)
bullets(s, Inches(1), Inches(3.6), Inches(9), Inches(2.5), [
    "PMs don't know if their strategy is the right one — execution and strategy are siloed",
    "It's laborious and slow to collate and synthesise insights across tools",
    "Decisions end up driven by opinion, politics, or the loudest voice",
    "Our customers leave us because they have to, not because they want to",
], sz=16, c=WHITE, bc=YELLOW, gap=Pt(10))


# ═══════════════════════════════════════════════════
# 4. WHAT CUSTOMERS ACTUALLY SAY
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.5),
    "What we're hearing", sz=14, b=True, c=PURPLE)

quotes = [
    ('"We have a general lack of data to support product decisions."', "Enterprise customer"),
    ('"Is Miro an appropriate tool for road mapping from ideation through delivery?"',
     "Anand Singh, Director IT, SMART Technologies"),
    ('"ProductBoard is deprecating the old ARR boards… we need to look at\ncompetitive offerings like Miro Insights."',
     "Enterprise customer evaluating alternatives"),
    ('"My target for this year was condensing our tech stack… I just consolidate\nall our project management tools from Asana, Airtable, Trello, Google Sheets…"',
     "AIHR Academy"),
]
y = Inches(1.2)
for q, src in quotes:
    rrect(s, Inches(1), y, Inches(10.5), Inches(1.15), WHITE)
    rect(s, Inches(1), y, Inches(0.05), Inches(1.15), PURPLE)
    txt(s, Inches(1.3), y + Inches(0.1), Inches(9.8), Inches(0.7), q, sz=13, c=TXT)
    txt(s, Inches(1.3), y + Inches(0.8), Inches(9.8), Inches(0.3), f"— {src}", sz=11, c=MUTED)
    y += Inches(1.35)


# ═══════════════════════════════════════════════════
# 5. WHY NOW (short)
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.8),
    "Why now", sz=34, b=True, c=WHITE)
bullets(s, Inches(1), Inches(1.8), Inches(10), Inches(4.5), [
    "Productboard is deprecating legacy boards — their customers are actively shopping",
    "We already have the building blocks: Timeline (WACU doubled), Tables, Kanban, Insights, Jira sync",
    "Teams don't want to leave Miro but feel forced to because of missing capabilities",
    "Estimated expansion ARR from roadmapping: $800K",
    "If we don't do this, someone else will",
], sz=17, c=WHITE, bc=YELLOW, gap=Pt(14))


# ═══════════════════════════════════════════════════
# 6. WHAT WE KNOW FROM CUSTOMERS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
    "What customers are telling us to build", sz=30, b=True, c=TXT)
txt(s, Inches(1), Inches(1.3), Inches(10), Inches(0.5),
    "From solution review, CSAT, deal rooms, and engagement data:", sz=14, c=MUTED)

rows = [
    ["Signal", "Data", "What it tells us"],
    ["44% say no tool gives them the viz they need", "Solution review", "We need an opinionated, ready-to-use roadmap space"],
    ["23% can't connect roadmap to company goals", "Solution review", "Roadmap needs strategic context, not just tasks"],
    ["200+ mentions of grouping & hierarchy", "CSAT feedback", "Table + Timeline + Kanban need to work together"],
    ["Jira-connected users have 4x active days", "Product analytics", "Jira integration is non-negotiable in V1"],
    ["Kanban has 9.1% multi-day engagement", "Product analytics", "Include Kanban as a core view"],
    ["$2.7M ARR in explicit roadmapping requests", "SFDC", "Named accounts: E.ON, Ubisoft, PWC, Keyloop, Fidelity"],
]
tbl(s, Inches(1), Inches(1.9), Inches(11.2), rows, [Inches(4.2), Inches(2), Inches(5)])


# ═══════════════════════════════════════════════════
# 7. WHAT V1 ACTUALLY IS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
    "V1 is three things", sz=34, b=True, c=WHITE)

# Col 1
rrect(s, Inches(1), Inches(1.7), Inches(3.5), Inches(4.8), DEEP)
rect(s, Inches(1), Inches(1.7), Inches(3.5), Inches(0.05), PURPLE)
txt(s, Inches(1.3), Inches(1.95), Inches(3), Inches(0.4),
    "A dedicated entry point", sz=18, b=True, c=WHITE)
txt(s, Inches(1.3), Inches(2.5), Inches(2.9), Inches(1.6),
    '"Product Roadmap" shows up in Create New and gets its own section '
    'in the spaces side panel. Roadmapping becomes a first-class thing in Miro, '
    'not something you hack together on a blank board.',
    sz=13, c=SOFT)
rect(s, Inches(1.3), Inches(4.3), Inches(2.9), Inches(0.01), RGBColor(0x30, 0x30, 0x60))
txt(s, Inches(1.3), Inches(4.5), Inches(2.9), Inches(0.3),
    "Solves: discoverability, getting started", sz=11, b=True, c=MUTED)

# Col 2
rrect(s, Inches(4.8), Inches(1.7), Inches(3.5), Inches(4.8), DEEP)
rect(s, Inches(4.8), Inches(1.7), Inches(3.5), Inches(0.05), GREEN)
txt(s, Inches(5.1), Inches(1.95), Inches(3), Inches(0.4),
    "A preconfigured space", sz=18, b=True, c=WHITE)
txt(s, Inches(5.1), Inches(2.5), Inches(2.9), Inches(1.6),
    'Three sections: a Pulse page with signals and suggested actions, '
    'a Backlog enriched with customer data, and a Roadmap with Table, '
    'Timeline, and Kanban views. Not a blank canvas — a starting point.',
    sz=13, c=SOFT)
rect(s, Inches(5.1), Inches(4.3), Inches(2.9), Inches(0.01), RGBColor(0x30, 0x30, 0x60))
txt(s, Inches(5.1), Inches(4.5), Inches(2.9), Inches(0.3),
    "Solves: 44% viz gap, view fragmentation", sz=11, b=True, c=MUTED)

# Col 3
rrect(s, Inches(8.6), Inches(1.7), Inches(3.5), Inches(4.8), DEEP)
rect(s, Inches(8.6), Inches(1.7), Inches(3.5), Inches(0.05), ORANGE)
txt(s, Inches(8.9), Inches(1.95), Inches(3), Inches(0.4),
    "Brownfield onboarding", sz=18, b=True, c=WHITE)
txt(s, Inches(8.9), Inches(2.5), Inches(2.9), Inches(1.6),
    'You already have a Jira backlog, Miro boards, strategy docs. '
    'We connect to those and populate your roadmap from what exists. '
    'AI does a first pass — you refine. Minutes, not hours. Even if manual at first.',
    sz=13, c=SOFT)
rect(s, Inches(8.9), Inches(4.3), Inches(2.9), Inches(0.01), RGBColor(0x30, 0x30, 0x60))
txt(s, Inches(8.9), Inches(4.5), Inches(2.9), Inches(0.3),
    "Solves: manual gathering, staleness", sz=11, b=True, c=MUTED)


# ═══════════════════════════════════════════════════
# 8. V0 → V1 — WHAT CHANGES
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
    "What we have vs. what we're building", sz=30, b=True, c=TXT)

# V0
rrect(s, Inches(1), Inches(1.6), Inches(4.8), Inches(4.8), WHITE)
txt(s, Inches(1.3), Inches(1.7), Inches(4), Inches(0.4), "Today", sz=16, b=True, c=MUTED)
bullets(s, Inches(1.3), Inches(2.2), Inches(4.2), Inches(3.8), [
    "Timeline — live, growing, but limited grouping and hierarchy",
    "Tables + Kanban — live, but disconnected from Timeline and Insights",
    "Miro Insights — live, but not connected to where PMs plan",
    "Jira sync — basic, status only",
    "AI Sidekick — on-demand, not proactive",
], sz=13, c=TXT, bc=MUTED, gap=Pt(8))
txt(s, Inches(1.3), Inches(5.3), Inches(4.2), Inches(0.4),
    "The pieces exist. They're just islands.", sz=12, b=True, c=RED)

# Arrow
txt(s, Inches(5.7), Inches(3.5), Inches(1), Inches(0.6), "→", sz=36, b=True, c=PURPLE, al=PP_ALIGN.CENTER)

# V1
sh = rrect(s, Inches(6.5), Inches(1.6), Inches(5.3), Inches(4.8), WHITE)
sh.line.color.rgb = PURPLE; sh.line.width = Pt(1.5)
txt(s, Inches(6.8), Inches(1.7), Inches(4), Inches(0.4), "V1", sz=16, b=True, c=PURPLE)
bullets(s, Inches(6.8), Inches(2.2), Inches(4.7), Inches(3.8), [
    'Entry point — "Product Roadmap" in Create New and spaces sidebar',
    "Preconfigured space — Pulse + Backlog + Roadmap, structured",
    "Brownfield — AI populates from Jira, boards, CSVs, docs",
    "Enrichment — customer count, ARR, quotes on every backlog item",
    "Pulse — surfaces what changed, what matters, what to do next",
], sz=13, c=TXT, bc=PURPLE, gap=Pt(8))
txt(s, Inches(6.8), Inches(5.3), Inches(4.7), Inches(0.4),
    "The PM makes decisions instead of gathering data.", sz=12, b=True, c=PURPLE)


# ═══════════════════════════════════════════════════
# 9. PAIN → SOLUTION MAP
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
    "How V1 maps to the actual pains", sz=30, b=True, c=WHITE)

rows = [
    ["What PMs struggle with", "Today", "V1"],
    ["Can't trace decisions to customer evidence",
     "Insights exists but lives elsewhere",
     "Enriched backlog with customer data. Side panel shows signals per item."],
    ["Roadmap goes stale within days",
     "Everything is manually maintained",
     "Jira sync keeps status current. Insights auto-refreshes. Pulse flags changes."],
    ["Spend 4-6 hrs/wk gathering signals",
     "Have to go find insights yourself",
     "Pulse surfaces what matters. Each signal has an action to add to backlog."],
    ["Can't validate ideas against real demand",
     "No link between feedback and backlog",
     "Every item shows mentions, ARR, companies. ICE scoring."],
    ["No tool for the visualisation (44%)",
     "Timeline, Table, Kanban are separate",
     "Unified views in one space. Same data, different lenses."],
]
tbl(s, Inches(1), Inches(1.6), Inches(11.2), rows,
    [Inches(3), Inches(3), Inches(5.2)], hbg=DEEP)


# ═══════════════════════════════════════════════════
# 10. TIMELINE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
    "Timeline", sz=30, b=True, c=TXT)

# POC
rrect(s, Inches(1), Inches(1.6), Inches(5), Inches(3.5), WHITE)
rect(s, Inches(1), Inches(1.6), Inches(5), Inches(0.05), ORANGE)
txt(s, Inches(1.3), Inches(1.85), Inches(4.4), Inches(0.4),
    "Smallest POC — mid March", sz=18, b=True, c=TXT)
bullets(s, Inches(1.3), Inches(2.4), Inches(4.4), Inches(2.2), [
    "Consolidated side panel in full screen (Miro + Jira + Insights data)",
    "Import existing backlog items from Jira / CSV",
    "Table enriched with Insights columns, refreshable",
], sz=14, c=TXT, bc=ORANGE, gap=Pt(10))

# MVP
rrect(s, Inches(6.5), Inches(1.6), Inches(5.5), Inches(3.5), WHITE)
rect(s, Inches(6.5), Inches(1.6), Inches(5.5), Inches(0.05), PURPLE)
txt(s, Inches(6.8), Inches(1.85), Inches(5), Inches(0.4),
    "Killer MVP — end of Q1", sz=18, b=True, c=TXT)
bullets(s, Inches(6.8), Inches(2.4), Inches(5), Inches(2.2), [
    "AI creation and population of Product Roadmap space",
    "Full-screen improvements for roadmap review and backlog grooming",
    "Timeline improvements",
    "Pulse page with signals and suggested actions",
], sz=14, c=TXT, bc=PURPLE, gap=Pt(10))

# Key questions
txt(s, Inches(1), Inches(5.5), Inches(10), Inches(0.4),
    "Open question: can we launch even the smallest POC without AI generation? Current thinking: no.",
    sz=13, c=MUTED)


# ═══════════════════════════════════════════════════
# 11. SUCCESS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.5),
    "How we'll know it's working", sz=30, b=True, c=WHITE)

rows = [
    ["Metric", "Target", "Why it matters"],
    ["Roadmaps created (90 days)", "500+", "Are people finding and using it?"],
    ["WAU per roadmap space", "5+", "Is the team actually living here?"],
    ["Jira connection rate", "70%+", "Connected users stick 4x more"],
    ["Insights connection rate", "50%+", "This is the differentiator"],
    ["Roadmap CSAT", "4.0+ (from 3.0)", "Are we actually less frustrating?"],
    ["Frustration score (12 mo)", "<30% (from 66%)", "The north star"],
]
tbl(s, Inches(1), Inches(1.3), Inches(11.2), rows,
    [Inches(3.2), Inches(2), Inches(6)], hbg=DEEP)


# ═══════════════════════════════════════════════════
# 12. RACI
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, BG)
txt(s, Inches(1), Inches(0.6), Inches(10), Inches(0.5),
    "RACI", sz=30, b=True, c=TXT)

rows = [
    ["Workstream", "Responsible", "Accountable", "Consulted", "Informed"],
    ["Product strategy", "Rosanna", "Rosanna", "Leader, Design", "Eng, GTM"],
    ["Entry point & space", "Eng (Platform)", "Eng Lead", "PM, Design", "Leader"],
    ["Brownfield & Pulse", "Eng (AI + FE)", "Eng Lead", "PM, AI, Design", "Leader"],
    ["Backlog & views", "Eng (Frontend)", "Eng Lead", "PM, Design", "Leader"],
    ["Jira sync", "Eng (Integrations)", "Eng Lead", "PM", "Leader"],
    ["Insights × Tables", "Lauren Brucato", "Lauren B.", "Rosanna, Eng", "Leader"],
    ["AI layer", "Eng (AI)", "AI Lead", "PM, Insights PM", "Leader"],
    ["GTM & launch", "Prod. Marketing", "GTM Lead", "PM, Leader", "Eng, CS"],
]
tbl(s, Inches(1), Inches(1.2), Inches(11.2), rows,
    [Inches(2.4), Inches(2), Inches(1.8), Inches(2.5), Inches(2.5)])

# Deps
txt(s, Inches(1), Inches(5.0), Inches(5), Inches(0.35),
    "Dependencies", sz=16, b=True, c=TXT)
dep_rows = [
    ["What", "Who", "Risk"],
    ["Insights × Tables enrichment GA", "Lauren Brucato", "Medium — must ship with V1"],
    ["Jira real-time sync", "Eng (Integrations)", "Low — foundation exists"],
    ["Spaces sidebar support", "Platform", "Low — exists"],
    ["AI inference infra", "AI Platform", "Low — powers Sidekick already"],
]
tbl(s, Inches(1), Inches(5.4), Inches(8), dep_rows, [Inches(3), Inches(2.5), Inches(2.5)])


# ═══════════════════════════════════════════════════
# 13. WHERE THIS GOES
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(1), Inches(0.8), Inches(10), Inches(0.8),
    "Where this goes", sz=30, b=True, c=WHITE)

txt(s, Inches(1), Inches(1.8), Inches(2.5), Inches(0.4), "V1 — now", sz=18, b=True, c=YELLOW)
txt(s, Inches(1), Inches(2.3), Inches(4), Inches(1.5),
    "Connect the pieces. Give PMs a home for roadmapping\n"
    "that's enriched by customer data and connected to execution.\n"
    "Dedicated entry point. Preconfigured space. Brownfield onboarding.",
    sz=14, c=SOFT)

txt(s, Inches(1), Inches(3.8), Inches(2.5), Inches(0.4), "V2 — next", sz=18, b=True, c=YELLOW)
txt(s, Inches(1), Inches(4.3), Inches(4), Inches(1),
    "Portfolio views for leaders. Goal alignment.\n"
    "Scenario planning. Automated status reporting.",
    sz=14, c=SOFT)

txt(s, Inches(1), Inches(5.3), Inches(2.5), Inches(0.4), "V3 — later", sz=18, b=True, c=YELLOW)
txt(s, Inches(1), Inches(5.8), Inches(4), Inches(1),
    "Roadmap autopilot. Agents that monitor signals,\n"
    "flag risks, propose changes. Humans steer, AI executes.",
    sz=14, c=SOFT)

# Right side — one line summary
txt(s, Inches(6.5), Inches(2.5), Inches(5.5), Inches(2),
    "The roadmap becomes a living\nsystem that gets smarter\nevery day, so the cost of\nstaying aligned drops to\nnear zero.",
    sz=24, b=True, c=WHITE)


# ═══════════════════════════════════════════════════
# 14. END
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BL); bg(s, DARK)
txt(s, Inches(0), Inches(2.8), W, Inches(1), "Let's go.", sz=56, b=True, c=WHITE, al=PP_ALIGN.CENTER)
rect(s, Inches(6.15), Inches(3.9), Inches(1), Inches(0.05), YELLOW)


# ── Save ──
out = '/Users/rosannaknottenbelt/Testing/Dynamic-Roadmap-V1-Kickoff.pptx'
prs.save(out)
print(f"Done → {out}")
